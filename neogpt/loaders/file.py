import json
import os
import tempfile
import zipfile

from loaders.base import BaseLoader


class TxtLoader(BaseLoader):
    def load(self):
        with open(self.path) as f:
            return f.read()

    def lazy_load(self):
        with open(self.path) as f:
            yield from f


class JSONLoader(BaseLoader):
    def load(self):
        with open(self.path) as f:
            return json.load(f)

    def lazy_load(self):
        with open(self.path) as f:
            yield from json.load(f)


class PDFLoader(BaseLoader):
    def __init__(self, path: str):
        try:
            from pdfminer.high_level import extract_text
        except ImportError as err:
            raise ImportError("""
                Please install pdfminer.six to use the PDFLoader.
                You can install it using:
                `pip install pdfminer.six`
            """) from err

        super().__init__(path)

    def load(self):
        from pdfminer.high_level import extract_text

        return extract_text(self.path).strip()

    def lazy_load(self):
        from pdfminer.high_level import extract_text

        yield extract_text(self.path)


class CSVLoader(BaseLoader):
    def __init__(self, path: str):
        try:
            import pandas as pd
        except ImportError as err:
            raise ImportError("""
                Please install pandas to use the CSVLoader.
                You can install it using:
                `pip install pandas`
            """) from err

        super().__init__(path)

    def load(self):
        import pandas as pd

        return pd.read_csv(self.path, engine="pyarrow")

    def lazy_load(self):
        import pandas as pd

        yield pd.read_csv(self.path)


class TSVLoader(BaseLoader):
    def __init__(self, path: str):
        try:
            import pandas as pd
        except ImportError as err:
            raise ImportError("""
                Please install pandas to use the TSVLoader.
                You can install it using:
                `pip install pandas`
            """) from err

        super().__init__(path)

    def load(self):
        import pandas as pd

        return pd.read_csv(self.path, sep="\t", engine="pyarrow")

    def lazy_load(self):
        import pandas as pd

        yield pd.read_csv(self.path, sep="\t")


class MarkdownLoader(BaseLoader):
    def __init__(self, path: str):
        try:
            import markdown
        except ImportError as err:
            raise ImportError("""
                Please install markdown to use the MarkdownLoader.
                You can install it using:
                `pip install markdown`
            """) from err

        super().__init__(path)

    def load(self):
        import markdown

        with open(self.path) as f:
            return markdown.markdown(f.read())

    def lazy_load(self):
        import markdown

        with open(self.path) as f:
            yield markdown.markdown(f.read())


class HTMLLoader(BaseLoader):
    def __init__(self, path: str):
        try:
            from bs4 import BeautifulSoup
        except ImportError as err:
            raise ImportError("""
                Please install beautifulsoup4 to use the HTMLLoader.
                You can install it using:
                `pip install beautifulsoup4`
            """) from err

        super().__init__(path)

    def load(self):
        from bs4 import BeautifulSoup

        with open(self.path, encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
            # Extract text content excluding HTML tags
            return soup.get_text(separator="\n", strip=True)

    def lazy_load(self):
        from bs4 import BeautifulSoup

        with open(self.path, encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
            # Yield clean text content in chunks
            yield soup.get_text(separator="\n", strip=True)


class DocxLoader(BaseLoader):
    def __init__(self, path: str):
        try:
            from docx import Document
        except ImportError as err:
            raise ImportError("""
                Please install python-docx to use the DocxLoader.
                You can install it using:
                `pip install python-docx`
            """) from err

        super().__init__(path)

    def load(self):
        from docx import Document

        doc = Document(self.path)
        return "\n".join([p.text for p in doc.paragraphs])

    def lazy_load(self):
        from docx import Document

        doc = Document(self.path)
        for p in doc.paragraphs:
            yield p.text


class PPTXLoader(BaseLoader):
    def __init__(self, path: str):
        try:
            from pptx import Presentation
        except ImportError as err:
            raise ImportError("""
                Please install python-pptx to use the PPTXLoader.
                You can install it using:
                `pip install python-pptx`
            """) from err

        super().__init__(path)

    def load(self):
        from pptx import Presentation

        prs = Presentation(self.path)
        return "\n".join([slide.text for slide in prs.slides])

    def lazy_load(self):
        from pptx import Presentation

        prs = Presentation(self.path)
        for slide in prs.slides:
            yield slide.text


class XLSXLoader(BaseLoader):
    def __init__(self, path: str):
        try:
            import pandas as pd
        except ImportError as err:
            raise ImportError("""
                Please install pandas to use the XLSXLoader.
                You can install it using:
                `pip install pandas`
            """) from err

        super().__init__(path)

    def load(self):
        import pandas as pd

        return pd.read_excel(self.path)

    def lazy_load(self):
        import pandas as pd

        yield pd.read_excel(self.path)


class EPUBLoader(BaseLoader):
    def __init__(self, path: str):
        try:
            from bs4 import BeautifulSoup
            from ebooklib import epub
        except ImportError as err:
            raise ImportError("""
                Please install EbookLib to use the EPUBLoader.
                You can install it using:
                `pip install EbookLib`
            """) from err

        super().__init__(path)

    def _clean_epub_text(self, text):
        from bs4 import BeautifulSoup

        formatted_text = ""
        soup = BeautifulSoup(text, "html.parser")
        all_text = soup.find_all(text=True)
        remove_tags = [
            "[document]",
            "noscript",
            "header",
            "html",
            "meta",
            "head",
            "input",
            "script",
        ]
        for t in all_text:
            if t.parent.name not in remove_tags:
                formatted_text += f"{t} "

        return formatted_text

    def load(self):
        from ebooklib import epub

        book = epub.read_epub(self.path)
        text = ""
        for item in book.get_items():
            text += self._clean_epub_text(item.get_body_content())
        return text

    def lazy_load(self):
        from ebooklib import epub

        book = epub.read_epub(self.path)
        for item in book.get_items():
            yield self._clean_epub_text(item.get_body_content())


class ZipLoader(BaseLoader):
    """Add support for loading .7z and other archive formats here."""

    def _extract_file(self, file, temp_dir):
        ext = os.path.splitext(file)[1].lower()
        loader_class = LOADER_MAP.get(ext)

        if loader_class:
            file_path = os.path.join(temp_dir, file)
            loader = loader_class(file_path)
            try:
                content = loader.load()
                return content.encode("utf-8", errors="ignore").decode("utf-8")
            except (UnicodeDecodeError, AttributeError) as e:
                print(f"Error processing file {file}: {e}")
                return ""

        return ""

    def load(self):
        extracted_data = ""
        with zipfile.ZipFile(self.path, "r") as z:  # noqa: SIM117
            with tempfile.TemporaryDirectory() as temp_dir:
                z.extractall(temp_dir)
                for file in z.namelist():
                    # Ignore system files like __MACOSX and .DS_Store
                    if (
                        file.startswith("__MACOSX/")
                        or file.endswith(".DS_Store")
                        or file.startswith(".")
                    ):
                        continue
                    extracted_data += self._extract_file(file, temp_dir)
        return extracted_data

    def lazy_load(self):
        with zipfile.ZipFile(self.path, "r") as z:  # noqa: SIM117
            with tempfile.TemporaryDirectory() as temp_dir:
                z.extractall(temp_dir)
                for file in z.namelist():
                    # Ignore system files . files that are internal hidden files
                    if (
                        file.startswith("__MACOSX/")
                        or file.endswith(".DS_Store")
                        or file.split(".")
                    ):
                        continue
                    ext = os.path.splitext(file)[1].lower()
                    loader_class = LOADER_MAP.get(ext)
                    if loader_class:
                        file_path = os.path.join(temp_dir, file)
                        loader = loader_class(file_path)
                        try:
                            for content in loader.lazy_load():
                                yield content.encode("utf-8", errors="ignore").decode(
                                    "utf-8"
                                )
                        except (UnicodeDecodeError, AttributeError) as e:
                            print(f"Error processing file {file}: {e}")


class AudioLoader(BaseLoader):
    """Add support to transcribe audio files such as .mp3, .wav, .flac, etc."""

    pass


class ImageLoader(BaseLoader):
    """Add OCR support for image files such as .jpg, .png, .bmp, etc. (Surya or Pytessearct) Later add support for multi-modal extraction"""

    pass


class SQLLoader(BaseLoader):
    """Add support for loading data from SQL databases using SQLAlchemy or something else"""

    pass


class WebLoader(BaseLoader):
    """Add support for loading web pages and scraping content using Beautiful Soup or Scrapy"""

    def __init__(self, path: str):
        try:
            import requests
            from bs4 import BeautifulSoup
        except ImportError as err:
            raise ImportError("""
                Please install beautifulsoup4 and requests to use the WebLoader.
                You can install them using:
                `pip install beautifulsoup4 requests`
            """) from err

        super().__init__(path)

    def _get_loader_for_url(self, url: str):
        """
        Get the appropriate loader class for the given URL based on URL_MAP.

        Args:
            url (str): The URL to check against URL_MAP.

        Returns:
            loader_class: The loader class corresponding to the URL's extension or domain.
        """
        for key, loader_class in URL_MAP.items():
            if key in url:
                return loader_class
        return None

    def load(self):
        import requests
        from bs4 import BeautifulSoup

        with open(self.path) as f:
            urls = f.readlines()

        content = ""
        for url in urls:
            url = url.strip()

            # Check if the URL is valid and accessible
            try:
                response = requests.get(url)
                response.raise_for_status()
            except requests.exceptions.RequestException as e:
                print(f"Error accessing URL {url}: {e}")
                continue

            # Get loader class for the URL
            loader_class = self._get_loader_for_url(url)

            if loader_class:
                # Create loader instance and load content
                loader = loader_class(url)
                try:
                    loaded_content = loader.load()
                    content += loaded_content + "\n"
                except Exception as e:
                    print(f"Error loading content from {url}: {e}")
                    continue
            else:
                print(f"No loader found for URL {url} using WebLoader")
                # Use BeautifulSoup to scrape content from the web page
                soup = BeautifulSoup(response.text, "html.parser")
                content += soup.get_text(separator="\n", strip=True) + "\n"

        return content.strip()


class YoutubeLoader(BaseLoader):
    """Add support for loading data from Youtube videos"""

    def __init__(self, url: str):
        try:
            import youtube_transcript_api as yt
        except ImportError as err:
            raise ImportError("""
                Please install youtube_transcript_api and pytube to use the YoutubeLoader.
                You can install them using:
                `pip install youtube_transcript_api`
            """) from err

        super().__init__(url)

    def _extract_video_id(self):
        import re

        return re.search(r"v=([^&]+)", self.path).group(1)

    def load(self):
        import youtube_transcript_api as yt

        video_id = self._extract_video_id()
        transcript = yt.YouTubeTranscriptApi.get_transcript(video_id)
        content = " ".join([line["text"] for line in transcript])
        return content

    def lazy_load(self):
        import youtube_transcript_api as yt

        video_id = self._extract_video_id()
        transcript = yt.YouTubeTranscriptApi.get_transcript(video_id)
        for line in transcript:
            yield line["text"]


URL_MAP = {
    "youtube.com": YoutubeLoader,
    "youtu.be": YoutubeLoader,
}

LOADER_MAP = {
    ".txt": TxtLoader,
    ".json": JSONLoader,
    ".pdf": PDFLoader,
    ".csv": CSVLoader,
    ".tsv": TSVLoader,
    ".md": MarkdownLoader,
    ".html": HTMLLoader,
    ".doc": DocxLoader,
    ".docx": DocxLoader,
    ".pptx": PPTXLoader,
    ".ppt": PPTXLoader,
    ".xlsx": XLSXLoader,
    ".xls": XLSXLoader,
    ".epub": EPUBLoader,
    ".zip": ZipLoader,
    ".lzma": ZipLoader,
    ".bz2": ZipLoader,
    ".gz": ZipLoader,
    ".tar": ZipLoader,
    ".tgz": ZipLoader,
    ".tar.gz": ZipLoader,
    ".url": WebLoader,
}
